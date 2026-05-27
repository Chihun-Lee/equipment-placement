"""
실험실 배치 시뮬레이터 m1 — 모바일 테스트 & App Store 배포 가이드 PPT 생성
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ─── 색상 팔레트 (다크 테마) ───
BG = RGBColor(0x1A, 0x1A, 0x1C)
BG_ELEV = RGBColor(0x26, 0x26, 0x2A)
FG = RGBColor(0xF5, 0xF5, 0xF7)
FG_DIM = RGBColor(0xA0, 0xA0, 0xA8)
ACCENT = RGBColor(0x4E, 0xC9, 0xB0)
WARN = RGBColor(0xFF, 0x6B, 0x4A)
BORDER = RGBColor(0x3A, 0x3A, 0x40)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]


def add_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    return bg


def add_text(slide, left, top, width, height, text, *, size=18, bold=False, color=FG, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='Apple SD Gothic Neo'):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split('\n') if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return tb


def add_rect(slide, left, top, width, height, fill=BG_ELEV, line=None, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    if radius:
        s.adjustments[0] = 0.12
    return s


def add_pill(slide, left, top, width, height, text, fill=ACCENT, fg=RGBColor(0x11, 0x11, 0x11)):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    s.adjustments[0] = 0.5
    tf = s.text_frame
    tf.margin_left = tf.margin_right = Emu(40000)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = fg
    run.font.name = 'Apple SD Gothic Neo'
    return s


def add_header(slide, page_num, total, section):
    # 상단 띠
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.45))
    bar.fill.solid(); bar.fill.fore_color.rgb = BG_ELEV
    bar.line.fill.background()
    add_text(slide, Inches(0.5), Inches(0.08), Inches(8), Inches(0.3),
             'Lab Placement m1 · 모바일 테스트 & App Store 배포 가이드',
             size=11, color=FG_DIM)
    add_text(slide, Inches(11.0), Inches(0.08), Inches(2.2), Inches(0.3),
             f'{section}  ·  {page_num}/{total}',
             size=11, color=FG_DIM, align=PP_ALIGN.RIGHT)


def add_title(slide, text, subtitle=None, top=Inches(0.7)):
    add_text(slide, Inches(0.6), top, Inches(12.2), Inches(0.7), text,
             size=30, bold=True, color=FG)
    if subtitle:
        add_text(slide, Inches(0.6), top + Inches(0.75), Inches(12.2), Inches(0.4),
                 subtitle, size=15, color=ACCENT)
    # 밑줄
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), top + Inches(1.25), Inches(1.0), Inches(0.06))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def bullet_list(slide, left, top, width, items, size=14, gap=0.38):
    for i, it in enumerate(items):
        y = top + Inches(i * gap)
        # 점
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, y + Inches(0.14), Inches(0.1), Inches(0.1))
        dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        add_text(slide, left + Inches(0.22), y, width - Inches(0.22), Inches(0.4),
                 it, size=size, color=FG)


def code_box(slide, left, top, width, height, code, size=11):
    box = add_rect(slide, left, top, width, height, fill=RGBColor(0x10, 0x10, 0x14), line=BORDER)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(120000); tf.margin_right = Emu(120000)
    tf.margin_top = Emu(80000); tf.margin_bottom = Emu(80000)
    lines = code.strip().split('\n')
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = ln
        run.font.size = Pt(size)
        run.font.name = 'Menlo'
        # 구문 강조 (단순)
        if ln.strip().startswith('#') or ln.strip().startswith('//'):
            run.font.color.rgb = FG_DIM
        elif any(ln.strip().startswith(kw) for kw in ['npm', 'npx', 'pod', 'bun', 'yarn', 'git', 'open', 'cd', 'python', 'ionic', 'xcodebuild']):
            run.font.color.rgb = ACCENT
        else:
            run.font.color.rgb = FG


# ───────────────────────────────────────────────────────────
# 슬라이드 1: 타이틀
# ───────────────────────────────────────────────────────────
def slide_title():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    # 그라데이션 느낌 액센트 블록
    block = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(0.6), Inches(2.4), Inches(0.6))
    block.fill.solid(); block.fill.fore_color.rgb = ACCENT
    block.line.fill.background()
    block.adjustments[0] = 0.5
    tf = block.text_frame
    tf.margin_left = tf.margin_right = Emu(100000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = 'LabPlace m1'
    r.font.size = Pt(18); r.font.bold = True
    r.font.color.rgb = RGBColor(0x11, 0x11, 0x11)
    r.font.name = 'Apple SD Gothic Neo'

    add_text(s, Inches(0.6), Inches(1.7), Inches(12.2), Inches(1.2),
             '실험실 배치 시뮬레이터',
             size=52, bold=True, color=FG)
    add_text(s, Inches(0.6), Inches(2.9), Inches(12.2), Inches(0.8),
             '모바일 최적화 테스트 & App Store 배포 가이드',
             size=24, color=ACCENT)

    # 부제 설명
    add_text(s, Inches(0.6), Inches(4.1), Inches(12.2), Inches(0.5),
             'PC에서 완벽 테스트 → Capacitor로 iOS 앱 빌드 → App Store 배포까지',
             size=16, color=FG_DIM)

    # 메타 정보
    meta = [
        ('📄 페이지', '총 22장'),
        ('🎯 대상', 'm1.html 모바일 버전'),
        ('🛠 도구', 'Capacitor · Xcode · TestFlight'),
        ('👤 작성', 'chihunlee@kims.re.kr'),
    ]
    for i, (k, v) in enumerate(meta):
        x = Inches(0.6 + (i % 2) * 6.4)
        y = Inches(5.2 + (i // 2) * 0.55)
        box = add_rect(s, x, y, Inches(6.0), Inches(0.45), fill=BG_ELEV)
        add_text(s, x + Inches(0.2), y + Inches(0.08), Inches(1.5), Inches(0.3),
                 k, size=12, color=FG_DIM)
        add_text(s, x + Inches(1.7), y + Inches(0.08), Inches(4.2), Inches(0.3),
                 v, size=12, bold=True, color=FG)

    add_text(s, Inches(0.6), Inches(6.8), Inches(12.2), Inches(0.4),
             '2026.04.12',
             size=12, color=FG_DIM)

slide_title()


# ───────────────────────────────────────────────────────────
# 슬라이드 2: 목차
# ───────────────────────────────────────────────────────────
def slide_toc():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, 2, 22, 'TOC')
    add_title(s, '목차', 'Table of Contents')

    sections = [
        ('1', 'PART 1. 프로젝트 구조 이해', '3-5',
         ['폴더 구조와 파일 역할', 'm1.html 모바일 최적화 요소', 'PWA 구성 요소']),
        ('2', 'PART 2. PC에서 테스트', '6-11',
         ['로컬 서버 띄우기', 'Chrome DevTools 모바일 에뮬레이션',
          '실제 iPhone에서 테스트 (같은 Wi-Fi)', 'ngrok으로 외부 공유 테스트']),
        ('3', 'PART 3. PWA 설치 테스트', '12-13',
         ['iOS Safari: 홈 화면 추가', 'Android Chrome: 설치 배너']),
        ('4', 'PART 4. App Store 배포', '14-20',
         ['Capacitor 설정', 'Xcode 프로젝트 빌드', 'TestFlight 베타',
          'App Store Connect 심사 제출']),
        ('5', 'PART 5. 체크리스트 & 문제해결', '21-22',
         ['배포 전 체크리스트', '자주 겪는 이슈']),
    ]
    y = Inches(2.0)
    for num, title, pages, subs in sections:
        # 번호 박스
        nb = add_rect(s, Inches(0.6), y, Inches(0.6), Inches(0.6), fill=ACCENT)
        nb_tf = nb.text_frame; nb_tf.margin_left = nb_tf.margin_right = Emu(0)
        nb_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = nb_tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num; r.font.size = Pt(18); r.font.bold = True
        r.font.color.rgb = RGBColor(0x11, 0x11, 0x11); r.font.name = 'Apple SD Gothic Neo'
        add_text(s, Inches(1.4), y - Inches(0.02), Inches(8.0), Inches(0.4),
                 title, size=16, bold=True, color=FG)
        add_text(s, Inches(1.4), y + Inches(0.3), Inches(8.0), Inches(0.3),
                 ' · '.join(subs), size=11, color=FG_DIM)
        add_text(s, Inches(11.8), y + Inches(0.1), Inches(1.3), Inches(0.4),
                 f'p. {pages}', size=12, color=ACCENT, align=PP_ALIGN.RIGHT)
        y += Inches(1.0)

slide_toc()


# ───────────────────────────────────────────────────────────
# PART 1. 프로젝트 구조
# ───────────────────────────────────────────────────────────
def slide_part_divider(num, title, subtitle, page, total):
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, page, total, f'PART {num}')
    # 큰 번호
    add_text(s, Inches(0.6), Inches(1.5), Inches(6), Inches(2.5),
             f'PART {num}', size=80, bold=True, color=RGBColor(0x30, 0x30, 0x38))
    add_text(s, Inches(0.6), Inches(3.8), Inches(12), Inches(1.2),
             title, size=44, bold=True, color=FG)
    add_text(s, Inches(0.6), Inches(4.9), Inches(12), Inches(0.6),
             subtitle, size=18, color=ACCENT)
    # 하단 큰 선
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(5.9), Inches(12.2), Inches(0.08))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

slide_part_divider(1, '프로젝트 구조 이해', '파일 구성 & 모바일 UX 설계', 3, 22)


def slide_folder_structure():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, 4, 22, 'PART 1')
    add_title(s, '폴더 구조', '~/Code/equiment_placement/')

    tree = """
equiment_placement/
├── index.html          ← 기존 데스크톱 버전 (그대로 둠)
├── m1.html             ← 🆕 모바일 최적화 버전
├── manifest.json       ← 🆕 PWA 설정 (앱 이름/아이콘/테마)
├── sw.js               ← 🆕 Service Worker (오프라인 캐시)
├── icon-192.png        ← 🆕 아이콘 192×192
├── icon-512.png        ← 🆕 아이콘 512×512
└── CLAUDE.md           ← 프로젝트 컨텍스트

GitHub Pages 배포 URL:
  · 데스크톱: https://chihun-lee.github.io/equipment-placement/
  · 모바일:   https://chihun-lee.github.io/equipment-placement/m1.html
"""
    code_box(s, Inches(0.6), Inches(2.1), Inches(8.5), Inches(4.5), tree, size=13)

    # 오른쪽 설명 카드
    card = add_rect(s, Inches(9.4), Inches(2.1), Inches(3.4), Inches(4.5))
    add_text(s, Inches(9.6), Inches(2.25), Inches(3.2), Inches(0.4),
             '핵심 포인트', size=14, bold=True, color=ACCENT)
    points = [
        '단일 파일 구조',
        '빌드 도구 없음',
        'CDN Three.js 사용',
        'localStorage 자동저장',
        'PWA로 앱처럼 설치',
    ]
    for i, pt in enumerate(points):
        add_text(s, Inches(9.6), Inches(2.8) + Inches(i * 0.5), Inches(3.2), Inches(0.4),
                 '✓  ' + pt, size=13, color=FG)

slide_folder_structure()


def slide_mobile_ux():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, 5, 22, 'PART 1')
    add_title(s, 'm1.html 모바일 UX 설계', '데스크톱과 무엇이 다른가')

    # 비교 표
    rows = [
        ('영역', '데스크톱 (index.html)', '모바일 (m1.html)'),
        ('레이아웃', '350px 사이드바 + 뷰포트', '풀스크린 뷰포트 + 바텀 시트'),
        ('입력 방식', '마우스 + 키보드', '터치 + 외장 키보드 옵션'),
        ('폰트 크기', '11~13px', '14~16px (iOS 자동 줌 방지)'),
        ('최소 터치 영역', '고려 안함', '44×44 px (Apple HIG)'),
        ('단축키', 'R/Del/Enter/Esc', '컨텍스트 바 버튼 + 외장키 대체'),
        ('되돌리기', '사이드바 버튼', '앱바 상시 버튼 + 비활성화 표시'),
        ('화면 전환', '사이드바 스크롤', '바텀 시트 push/pop 스택'),
        ('뒤로가기', '없음', '앱바 ← 버튼 + popstate 핸들러'),
        ('선택 UI', '사이드바 편집 패널', '상단 컨텍스트 바 (회전/복제/편집/삭제)'),
        ('카메라', '사이드바 프리셋', '앱바 📷 버튼 → 카메라 시트'),
        ('설치', 'N/A', 'PWA manifest + 홈 화면 추가 프롬프트'),
    ]
    col_w = [Inches(2.0), Inches(4.8), Inches(5.4)]
    x0 = Inches(0.6); y0 = Inches(2.0)
    row_h = Inches(0.42)
    for ri, row in enumerate(rows):
        is_head = ri == 0
        y = y0 + row_h * ri
        bg_color = BG_ELEV if is_head else (RGBColor(0x20, 0x20, 0x24) if ri % 2 else BG)
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x0, y, sum(col_w, Emu(0)), row_h)
        bg.fill.solid(); bg.fill.fore_color.rgb = bg_color
        bg.line.color.rgb = BORDER; bg.line.width = Pt(0.5)
        x = x0
        for ci, cell in enumerate(row):
            add_text(s, x + Inches(0.15), y + Inches(0.09), col_w[ci] - Inches(0.3), row_h - Inches(0.1),
                     cell, size=11 if not is_head else 12,
                     bold=is_head, color=ACCENT if is_head else FG)
            x += col_w[ci]

slide_mobile_ux()


# ───────────────────────────────────────────────────────────
# PART 2. PC에서 테스트
# ───────────────────────────────────────────────────────────
slide_part_divider(2, 'PC에서 테스트', 'Local Server · DevTools · Real Device · ngrok', 6, 22)


def slide_why_server():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, 7, 22, 'PART 2')
    add_title(s, '왜 로컬 서버가 필요한가', 'file:// 로는 안 되는 이유')

    # 좌: 문제
    card_a = add_rect(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(4.8), fill=RGBColor(0x2A, 0x1E, 0x1E), line=WARN)
    add_text(s, Inches(0.8), Inches(2.1), Inches(5.6), Inches(0.4),
             '❌ file:// 로 열면', size=16, bold=True, color=WARN)
    probs = [
        'ES modules (importmap) 동작 안 함',
        'Service Worker 등록 불가',
        'manifest.json 로드 실패',
        'CORS 정책으로 fetch 차단',
        '실제 모바일과 동작이 달라짐',
    ]
    for i, p in enumerate(probs):
        add_text(s, Inches(0.8), Inches(2.8) + Inches(i * 0.5), Inches(5.6), Inches(0.4),
                 '• ' + p, size=13, color=FG)

    # 우: 해결
    card_b = add_rect(s, Inches(6.8), Inches(2.0), Inches(6.0), Inches(4.8), fill=RGBColor(0x1A, 0x2A, 0x26), line=ACCENT)
    add_text(s, Inches(7.0), Inches(2.1), Inches(5.6), Inches(0.4),
             '✅ 로컬 HTTP 서버로 열기', size=16, bold=True, color=ACCENT)
    sols = [
        'Python 내장 서버 (가장 간단)',
        'Node http-server (npm -g)',
        'VS Code Live Server 확장',
        'Bun serve · Caddy 등',
    ]
    for i, p in enumerate(sols):
        add_text(s, Inches(7.0), Inches(2.8) + Inches(i * 0.5), Inches(5.6), Inches(0.4),
                 '• ' + p, size=13, color=FG)

    # 추천 명령어
    add_text(s, Inches(7.0), Inches(5.1), Inches(5.6), Inches(0.4),
             '👉 이 가이드는 Python 서버를 사용합니다',
             size=12, color=FG_DIM)
    code_box(s, Inches(7.0), Inches(5.5), Inches(5.6), Inches(1.1),
             'cd ~/Code/equiment_placement\npython3 -m http.server 8000', size=12)

slide_why_server()


def slide_local_server():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, 8, 22, 'PART 2')
    add_title(s, '로컬 서버 띄우기', '3가지 옵션')

    options = [
        ('A', 'Python 내장 서버 (추천)', 'macOS 기본 설치',
         'cd ~/Code/equiment_placement\npython3 -m http.server 8000\n# → http://localhost:8000/m1.html'),
        ('B', 'Node http-server', 'npm 글로벌 설치 필요',
         'npm install -g http-server\ncd ~/Code/equiment_placement\nhttp-server -p 8000\n# → http://localhost:8000/m1.html'),
        ('C', 'VS Code Live Server', 'VS Code 확장 (GUI)',
         '# 1) VS Code 확장에서 "Live Server" 설치\n# 2) m1.html 우클릭 → Open with Live Server\n# → 자동으로 브라우저 열림'),
    ]
    y = Inches(1.95)
    for letter, title, desc, code in options:
        card = add_rect(s, Inches(0.6), y, Inches(12.2), Inches(1.55))
        # 글자 원
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), y + Inches(0.35), Inches(0.85), Inches(0.85))
        circle.fill.solid(); circle.fill.fore_color.rgb = ACCENT
        circle.line.fill.background()
        tf = circle.text_frame; tf.margin_left = tf.margin_right = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = letter; r.font.size = Pt(22); r.font.bold = True
        r.font.color.rgb = RGBColor(0x11, 0x11, 0x11); r.font.name = 'Apple SD Gothic Neo'
        add_text(s, Inches(1.9), y + Inches(0.12), Inches(6), Inches(0.4),
                 title, size=15, bold=True, color=FG)
        add_text(s, Inches(1.9), y + Inches(0.5), Inches(6), Inches(0.4),
                 desc, size=11, color=FG_DIM)
        code_box(s, Inches(6.9), y + Inches(0.15), Inches(5.8), Inches(1.25), code, size=10)
        y += Inches(1.7)

slide_local_server()


def slide_devtools():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, 9, 22, 'PART 2')
    add_title(s, 'Chrome DevTools 모바일 에뮬레이션', '빠른 1차 테스트')

    steps = [
        ('1', '브라우저에서 접속', 'http://localhost:8000/m1.html'),
        ('2', 'DevTools 열기', 'Cmd + Option + I (macOS)'),
        ('3', '디바이스 툴바', 'Cmd + Shift + M 또는 상단 📱 아이콘'),
        ('4', '디바이스 선택', '"iPhone 15 Pro", "Pixel 8" 등 프리셋 선택'),
        ('5', 'DPR 조정', '상단 바에서 DPR, 네트워크, 터치 모드 설정'),
        ('6', '회전 테스트', '상단 회전 아이콘으로 가로/세로 전환'),
    ]
    for i, (n, title, detail) in enumerate(steps):
        y = Inches(2.0 + i * 0.7)
        nb = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y + Inches(0.05), Inches(0.5), Inches(0.5))
        nb.fill.solid(); nb.fill.fore_color.rgb = ACCENT
        nb.line.fill.background()
        tf = nb.text_frame; tf.margin_left = tf.margin_right = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = n; r.font.size = Pt(14); r.font.bold = True
        r.font.color.rgb = RGBColor(0x11, 0x11, 0x11); r.font.name = 'Apple SD Gothic Neo'
        add_text(s, Inches(1.4), y + Inches(0.05), Inches(3.5), Inches(0.4),
                 title, size=14, bold=True, color=FG)
        add_text(s, Inches(4.9), y + Inches(0.1), Inches(8), Inches(0.4),
                 detail, size=12, color=FG_DIM, font='Menlo')

    # 주의 박스
    warn_box = add_rect(s, Inches(0.6), Inches(6.4), Inches(12.2), Inches(0.7), fill=RGBColor(0x2A, 0x24, 0x14), line=WARN)
    add_text(s, Inches(0.85), Inches(6.55), Inches(11.8), Inches(0.4),
             '⚠ 에뮬레이션은 레이아웃 확인용. 실제 터치 성능, 3D GPU 렌더링 속도는 반드시 실기기에서 검증하세요.',
             size=12, color=FG)

slide_devtools()


def slide_real_device():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, 10, 22, 'PART 2')
    add_title(s, '실제 iPhone에서 테스트', '같은 Wi-Fi 네트워크 사용')

    # 좌: 단계
    add_text(s, Inches(0.6), Inches(2.0), Inches(6.5), Inches(0.4),
             '단계별 절차', size=16, bold=True, color=ACCENT)

    steps_code = """
# 1) 맥북의 로컬 IP 확인
ifconfig | grep "inet " | grep -v 127.0.0.1
# → inet 192.168.0.42 ... 이런 줄에서 IP 복사

# 2) 서버를 0.0.0.0 으로 띄움 (모든 인터페이스)
cd ~/Code/equiment_placement
python3 -m http.server 8000 --bind 0.0.0.0

# 3) 맥북 방화벽에서 해당 포트 허용
#    시스템 설정 → 네트워크 → 방화벽 → 옵션 → Python 허용
"""
    code_box(s, Inches(0.6), Inches(2.45), Inches(6.5), Inches(3.4), steps_code, size=11)

    add_text(s, Inches(0.6), Inches(5.95), Inches(6.5), Inches(0.4),
             '📱 iPhone Safari에서 접속', size=14, bold=True, color=ACCENT)
    add_text(s, Inches(0.6), Inches(6.3), Inches(6.5), Inches(0.5),
             'http://192.168.0.42:8000/m1.html',
             size=14, color=FG, font='Menlo')
    add_text(s, Inches(0.6), Inches(6.65), Inches(6.5), Inches(0.5),
             '※ 맥북과 같은 Wi-Fi여야 함',
             size=11, color=FG_DIM)

    # 우: Safari Remote Debugging
    add_text(s, Inches(7.4), Inches(2.0), Inches(5.5), Inches(0.4),
             'Safari Web Inspector 연결', size=16, bold=True, color=ACCENT)
    steps = [
        ('1', 'iPhone: 설정 → Safari → 고급 → 웹 속성 검사 ON'),
        ('2', 'Mac: Safari → 환경설정 → 고급 → "개발자용 메뉴 표시" ON'),
        ('3', 'USB 케이블로 iPhone 연결 (Finder에서 신뢰 허용)'),
        ('4', 'Mac Safari → 개발자 → [iPhone 이름] → m1.html 선택'),
        ('5', 'JS 콘솔/네트워크/레이아웃 인스펙터를 iPhone 화면에 직접 연결'),
    ]
    for i, (n, t) in enumerate(steps):
        y = Inches(2.5 + i * 0.55)
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.4), y + Inches(0.05), Inches(0.4), Inches(0.4))
        circle.fill.solid(); circle.fill.fore_color.rgb = BG_ELEV
        circle.line.color.rgb = ACCENT; circle.line.width = Pt(1)
        tf = circle.text_frame; tf.margin_left = tf.margin_right = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = n; r.font.size = Pt(11); r.font.bold = True
        r.font.color.rgb = ACCENT; r.font.name = 'Apple SD Gothic Neo'
        add_text(s, Inches(7.9), y + Inches(0.08), Inches(5.2), Inches(0.5),
                 t, size=11, color=FG)

slide_real_device()


def slide_ngrok():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, 11, 22, 'PART 2')
    add_title(s, 'ngrok으로 외부 공유 테스트', '다른 네트워크의 사용자도 테스트 가능')

    add_text(s, Inches(0.6), Inches(2.0), Inches(12), Inches(0.4),
             '왜 ngrok? → HTTPS가 필요한 PWA/Service Worker/카메라 API 테스트 + 외부 공유',
             size=13, color=FG_DIM)

    code = """
# 1) ngrok 설치 (Homebrew)
brew install ngrok

# 2) 계정 가입 후 토큰 등록 (최초 1회)
ngrok config add-authtoken YOUR_TOKEN_HERE

# 3) 로컬 서버 실행
python3 -m http.server 8000

# 4) 별도 터미널에서 ngrok 실행
ngrok http 8000

# 출력 예:
#   Forwarding  https://abcd-1234.ngrok-free.app → http://localhost:8000
#
# → iPhone Safari에서 https://abcd-1234.ngrok-free.app/m1.html 접속
#   HTTPS이므로 Service Worker, Clipboard, 카메라, 홈화면 추가 전부 동작
"""
    code_box(s, Inches(0.6), Inches(2.5), Inches(8.5), Inches(4.5), code, size=11)

    # 우: 장단점
    card = add_rect(s, Inches(9.4), Inches(2.5), Inches(3.4), Inches(4.5))
    add_text(s, Inches(9.6), Inches(2.65), Inches(3.1), Inches(0.4),
             '장점', size=13, bold=True, color=ACCENT)
    pros = ['HTTPS 자동 제공', '외부에서 접속 가능', 'PWA 전체 기능 테스트', '방화벽 설정 불필요']
    for i, p in enumerate(pros):
        add_text(s, Inches(9.6), Inches(3.0) + Inches(i * 0.35), Inches(3.1), Inches(0.3),
                 '+ ' + p, size=11, color=FG)

    add_text(s, Inches(9.6), Inches(4.55), Inches(3.1), Inches(0.4),
             '단점', size=13, bold=True, color=WARN)
    cons = ['무료는 URL 변경됨', '세션 2시간 제한', '느릴 수 있음']
    for i, p in enumerate(cons):
        add_text(s, Inches(9.6), Inches(4.9) + Inches(i * 0.35), Inches(3.1), Inches(0.3),
                 '- ' + p, size=11, color=FG)

slide_ngrok()


# ───────────────────────────────────────────────────────────
# PART 3. PWA 설치 테스트
# ───────────────────────────────────────────────────────────
slide_part_divider(3, 'PWA 설치 테스트', '홈 화면에 추가하여 앱처럼 사용', 12, 22)


def slide_pwa_install():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, 13, 22, 'PART 3')
    add_title(s, 'iOS / Android 설치 방법', 'App Store 없이도 앱처럼 쓸 수 있음')

    # iOS
    ios_card = add_rect(s, Inches(0.6), Inches(2.0), Inches(6.0), Inches(5.0))
    add_text(s, Inches(0.8), Inches(2.15), Inches(5.6), Inches(0.4),
             '🍎 iOS Safari', size=18, bold=True, color=FG)
    add_text(s, Inches(0.8), Inches(2.5), Inches(5.6), Inches(0.4),
             '※ Chrome 아님, 반드시 Safari', size=11, color=FG_DIM)
    ios_steps = [
        ('1', 'Safari에서 배포 URL 열기'),
        ('2', '하단 공유 버튼 ⎋ 탭'),
        ('3', '"홈 화면에 추가" 선택'),
        ('4', '이름 확인 → "추가"'),
        ('5', '홈화면에 아이콘 생성됨'),
        ('6', '탭하면 풀스크린 앱처럼 실행'),
    ]
    for i, (n, t) in enumerate(ios_steps):
        y = Inches(2.95 + i * 0.6)
        add_text(s, Inches(0.8), y, Inches(0.4), Inches(0.4),
                 n + '.', size=13, bold=True, color=ACCENT)
        add_text(s, Inches(1.15), y, Inches(5.2), Inches(0.4),
                 t, size=13, color=FG)

    # Android
    and_card = add_rect(s, Inches(6.8), Inches(2.0), Inches(6.0), Inches(5.0))
    add_text(s, Inches(7.0), Inches(2.15), Inches(5.6), Inches(0.4),
             '🤖 Android Chrome', size=18, bold=True, color=FG)
    add_text(s, Inches(7.0), Inches(2.5), Inches(5.6), Inches(0.4),
             '자동 설치 배너 표시됨', size=11, color=FG_DIM)
    and_steps = [
        ('1', 'Chrome에서 배포 URL 열기'),
        ('2', '하단 "홈 화면에 추가" 배너 탭'),
        ('3', '또는 우상단 ⋮ → "앱 설치"'),
        ('4', '"설치" 클릭'),
        ('5', '앱 서랍에 아이콘 생성'),
        ('6', '주소창 없는 풀스크린 모드'),
    ]
    for i, (n, t) in enumerate(and_steps):
        y = Inches(2.95 + i * 0.6)
        add_text(s, Inches(7.0), y, Inches(0.4), Inches(0.4),
                 n + '.', size=13, bold=True, color=ACCENT)
        add_text(s, Inches(7.35), y, Inches(5.2), Inches(0.4),
                 t, size=13, color=FG)

    # 하단 안내
    add_text(s, Inches(0.6), Inches(7.05), Inches(12.2), Inches(0.4),
             '💡 PWA로도 충분하다면 App Store 배포를 건너뛸 수 있습니다. 진짜 앱 배포가 필요한 이유가 있을 때만 PART 4로.',
             size=12, color=FG_DIM)

slide_pwa_install()


# ───────────────────────────────────────────────────────────
# PART 4. App Store 배포
# ───────────────────────────────────────────────────────────
slide_part_divider(4, 'App Store 배포', 'Capacitor → Xcode → TestFlight → App Store', 14, 22)


def slide_capacitor_overview():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, 15, 22, 'PART 4')
    add_title(s, '왜 Capacitor인가', '웹뷰 기반 네이티브 래퍼 비교')

    # 비교 카드 3개
    options_data = [
        ('Capacitor', '추천', ACCENT,
         ['Ionic 팀 최신 유지', 'Xcode/Gradle 직접 제어', '플러그인 생태계 풍부', '문서화 우수'],
         'npm i @capacitor/core\nnpx cap add ios'),
        ('Cordova', '레거시', RGBColor(0x88, 0x88, 0x88),
         ['오래된 방식', '플러그인 호환성 이슈', '신규 프로젝트 권장 X', ''],
         '# 권장하지 않음'),
        ('WKWebView 직접', '고급자용', RGBColor(0xff, 0xaa, 0x44),
         ['Swift 코드 직접 작성', '빌드 파이프라인 수동', '자유도 최고', '학습 곡선 급함'],
         '// Xcode에서 iOS 앱 생성 후\n// WKWebView로 m1.html 로드'),
    ]
    for i, (name, badge, color, features, code) in enumerate(options_data):
        x = Inches(0.6 + i * 4.25)
        card = add_rect(s, x, Inches(2.0), Inches(4.1), Inches(4.9))
        # 상단 색 띠
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.0), Inches(4.1), Inches(0.12))
        bar.fill.solid(); bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_text(s, x + Inches(0.25), Inches(2.25), Inches(3.6), Inches(0.5),
                 name, size=20, bold=True, color=FG)
        add_pill(s, x + Inches(2.9), Inches(2.3), Inches(1.0), Inches(0.35), badge, fill=color)
        for j, f in enumerate(features):
            if not f: continue
            add_text(s, x + Inches(0.25), Inches(2.95) + Inches(j * 0.4), Inches(3.7), Inches(0.4),
                     '✓ ' + f, size=12, color=FG)
        code_box(s, x + Inches(0.25), Inches(4.7), Inches(3.7), Inches(2.0), code, size=10)

    add_text(s, Inches(0.6), Inches(7.1), Inches(12.2), Inches(0.4),
             '👉 이 가이드는 Capacitor로 진행합니다. m1.html의 코드 변경 없이 iOS/Android 네이티브 래퍼 생성 가능.',
             size=12, color=ACCENT)

slide_capacitor_overview()


def slide_capacitor_setup():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, 16, 22, 'PART 4')
    add_title(s, 'Capacitor 프로젝트 설정', '새 폴더에서 시작 (원본 건드리지 않음)')

    code = """
# 1) 새 폴더 만들고 Node 프로젝트 초기화
mkdir -p ~/Code/equiment_placement_ios
cd ~/Code/equiment_placement_ios
npm init -y

# 2) Capacitor 설치
npm install @capacitor/core @capacitor/cli
npm install @capacitor/ios

# 3) www 폴더에 m1.html + 에셋 복사
mkdir www
cp ~/Code/equiment_placement/m1.html www/index.html
cp ~/Code/equiment_placement/manifest.json www/
cp ~/Code/equiment_placement/sw.js www/
cp ~/Code/equiment_placement/icon-*.png www/

# 4) Capacitor 초기화
npx cap init "LabPlace" "re.kims.labplace" --web-dir=www

# 5) iOS 플랫폼 추가
npx cap add ios

# 6) 웹 코드 동기화 (수정 후 매번 실행)
npx cap sync ios

# 7) Xcode로 열기
npx cap open ios
"""
    code_box(s, Inches(0.6), Inches(2.0), Inches(8.5), Inches(5.2), code, size=11)

    # 오른쪽 설명
    add_text(s, Inches(9.4), Inches(2.0), Inches(3.4), Inches(0.4),
             '📌 핵심 파라미터', size=14, bold=True, color=ACCENT)
    params = [
        ('App 이름', 'LabPlace'),
        ('Bundle ID', 're.kims.labplace'),
        ('웹 디렉토리', 'www/'),
        ('진입 파일', 'index.html'),
        ('플랫폼', 'ios'),
    ]
    for i, (k, v) in enumerate(params):
        y = Inches(2.55 + i * 0.45)
        add_text(s, Inches(9.4), y, Inches(1.5), Inches(0.4),
                 k, size=11, color=FG_DIM)
        add_text(s, Inches(10.7), y, Inches(2.2), Inches(0.4),
                 v, size=11, color=FG, font='Menlo', bold=True)

    note = add_rect(s, Inches(9.4), Inches(5.0), Inches(3.4), Inches(2.2),
                    fill=RGBColor(0x1A, 0x2A, 0x2E), line=ACCENT)
    add_text(s, Inches(9.55), Inches(5.1), Inches(3.2), Inches(0.4),
             '💡 Bundle ID', size=12, bold=True, color=ACCENT)
    add_text(s, Inches(9.55), Inches(5.45), Inches(3.2), Inches(1.7),
             '역도메인 형식.\n앱스토어에서 고유해야 함.\n한번 등록하면 변경 불가.\n소속 도메인 있으면 그것 추천.',
             size=11, color=FG)

slide_capacitor_setup()


def slide_xcode_build():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, 17, 22, 'PART 4')
    add_title(s, 'Xcode에서 빌드 & 실행', '시뮬레이터 → 실기기 순서')

    steps = [
        ('1', 'Xcode 프로젝트 열림', 'npx cap open ios 실행 후 자동으로 Xcode 시작'),
        ('2', 'Signing & Capabilities 탭', 'Team을 본인 Apple ID로 설정\n(무료 계정도 7일 한정으로 실기기 설치 가능)'),
        ('3', 'Bundle Identifier 확인', 're.kims.labplace\n(Apple Developer 계정에서 App ID 등록 필요)'),
        ('4', 'Info.plist 수정', 'Status bar, 지원 방향, 앱 이름 표시 설정\n설명이 필요한 권한은 NS...UsageDescription 추가'),
        ('5', '시뮬레이터 실행', '상단 기기 선택 → iPhone 15 Pro → ▶ Run\n첫 빌드는 2~5분 소요'),
        ('6', '실기기 실행', 'iPhone USB 연결 → 기기 선택 → ▶\n설정 → 일반 → 기기관리에서 개발자 신뢰'),
    ]
    for i, (n, title, desc) in enumerate(steps):
        y = Inches(2.0 + i * 0.85)
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y + Inches(0.15), Inches(0.55), Inches(0.55))
        circle.fill.solid(); circle.fill.fore_color.rgb = ACCENT
        circle.line.fill.background()
        tf = circle.text_frame; tf.margin_left = tf.margin_right = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = n; r.font.size = Pt(14); r.font.bold = True
        r.font.color.rgb = RGBColor(0x11, 0x11, 0x11); r.font.name = 'Apple SD Gothic Neo'
        add_text(s, Inches(1.35), y + Inches(0.05), Inches(4), Inches(0.5),
                 title, size=14, bold=True, color=FG)
        add_text(s, Inches(5.4), y + Inches(0.05), Inches(7.5), Inches(0.8),
                 desc, size=11, color=FG_DIM)

slide_xcode_build()


def slide_developer_account():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, 18, 22, 'PART 4')
    add_title(s, 'Apple Developer Program', '연 $99 · 가입 없이는 App Store 배포 불가')

    # 단계별 설명
    phases = [
        ('무료 Apple ID', '$0',
         '• 시뮬레이터 빌드 ✓\n• 실기기 7일 임시 설치 ✓\n• App Store 배포 ✗\n• TestFlight ✗'),
        ('Developer Program 개인', '$99/년',
         '• 실기기 영구 설치 ✓\n• App Store 배포 ✓\n• TestFlight (최대 10,000명) ✓\n• 푸시/인앱결제 ✓'),
        ('Developer Program 조직', '$99/년',
         '• 위의 모든 혜택 +\n• 팀 멤버 권한 관리\n• D-U-N-S 번호 필요\n• 법인 등록증 필요'),
    ]
    for i, (title, price, desc) in enumerate(phases):
        x = Inches(0.6 + i * 4.25)
        card = add_rect(s, x, Inches(2.0), Inches(4.1), Inches(4.0))
        add_text(s, x + Inches(0.25), Inches(2.15), Inches(3.7), Inches(0.5),
                 title, size=16, bold=True, color=FG)
        add_pill(s, x + Inches(0.25), Inches(2.65), Inches(1.2), Inches(0.35), price,
                 fill=ACCENT if i == 1 else BG_ELEV,
                 fg=RGBColor(0x11, 0x11, 0x11) if i == 1 else FG)
        add_text(s, x + Inches(0.25), Inches(3.2), Inches(3.7), Inches(2.7),
                 desc, size=12, color=FG)

    # 가입 절차
    add_text(s, Inches(0.6), Inches(6.3), Inches(12), Inches(0.4),
             '가입 절차', size=14, bold=True, color=ACCENT)
    code_box(s, Inches(0.6), Inches(6.65), Inches(12.2), Inches(0.6),
             '# https://developer.apple.com/programs/enroll/ → 로그인 → 개인/법인 선택 → 결제 → 24~48시간 승인',
             size=10)

slide_developer_account()


def slide_testflight():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, 19, 22, 'PART 4')
    add_title(s, 'TestFlight로 베타 배포', '심사 전 내부/외부 테스터에게 배포')

    steps = [
        '1. Xcode → Product → Archive (빌드 아카이브 생성)',
        '2. Organizer 창이 열림 → Distribute App → App Store Connect',
        '3. Upload 선택 → 인증서/프로파일 자동 관리 ✓',
        '4. 업로드 완료 (5~15분) → App Store Connect에서 처리 대기',
        '5. appstoreconnect.apple.com → 앱 선택 → TestFlight 탭',
        '6. "내부 테스트 그룹" 생성 → 본인 계정 추가 → 빌드 활성화',
        '7. iPhone에 TestFlight 앱 설치 → 초대 수락 → 앱 설치',
        '8. 외부 테스터는 심사(1~2일) 필요. 내부 테스터는 즉시 설치 가능',
    ]
    for i, step in enumerate(steps):
        y = Inches(2.0 + i * 0.55)
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y + Inches(0.05), Inches(0.45), Inches(0.45))
        circle.fill.solid(); circle.fill.fore_color.rgb = BG_ELEV
        circle.line.color.rgb = ACCENT; circle.line.width = Pt(1.5)
        add_text(s, Inches(1.2), y + Inches(0.08), Inches(12), Inches(0.4),
                 step, size=13, color=FG)

    # 하단 팁
    tip = add_rect(s, Inches(0.6), Inches(6.5), Inches(12.2), Inches(0.8),
                   fill=RGBColor(0x1A, 0x2A, 0x2E), line=ACCENT)
    add_text(s, Inches(0.8), Inches(6.65), Inches(11.8), Inches(0.5),
             '💡 TestFlight 빌드는 최대 90일 유효. 만료 전 새 빌드 업로드 필요.',
             size=12, color=FG)

slide_testflight()


def slide_app_store_submit():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, 20, 22, 'PART 4')
    add_title(s, 'App Store 심사 제출', '최종 배포까지 남은 단계')

    # 좌: 준비물
    add_text(s, Inches(0.6), Inches(2.0), Inches(6), Inches(0.4),
             '📋 제출 전 필수 준비물', size=15, bold=True, color=ACCENT)
    items = [
        '앱 아이콘 1024×1024 PNG (알파 없음)',
        '스크린샷: 6.7"(iPhone 15 Pro Max) 최소 3장',
        '스크린샷: 6.5" · 5.5" 호환 기기별',
        '앱 이름 (30자 이내)',
        '부제 (30자 이내)',
        '설명 (4000자 이내)',
        '키워드 (100자 이내, 쉼표 구분)',
        '지원 URL / 마케팅 URL',
        '프라이버시 정책 URL (웹에 호스팅)',
        '카테고리 (1차/2차)',
        '연령 등급 설문 응답',
        '데이터 수집 내역 (Privacy Nutrition Label)',
    ]
    for i, it in enumerate(items):
        y = Inches(2.5 + (i // 2) * 0.45)
        x = Inches(0.6 + (i % 2) * 3.0)
        add_text(s, x, y, Inches(3), Inches(0.4),
                 '□ ' + it, size=11, color=FG)

    # 우: 절차
    add_text(s, Inches(7.0), Inches(2.0), Inches(6), Inches(0.4),
             '🚀 제출 절차', size=15, bold=True, color=ACCENT)
    steps = [
        'App Store Connect → 내 앱 → "+ 신규"',
        '플랫폼 iOS · 이름 · Bundle ID 입력',
        '앱 정보 입력 (위 준비물 모두)',
        '스크린샷 업로드',
        'TestFlight에서 사용할 빌드 선택',
        '앱 심사 노트에 테스트 계정 정보',
        '"심사를 위해 제출" 클릭',
        '심사 대기 (평균 24~48시간)',
        '승인 후 "수동 출시" 또는 "즉시 출시"',
        '🎉 App Store에 공개',
    ]
    for i, step in enumerate(steps):
        y = Inches(2.5 + i * 0.42)
        add_text(s, Inches(7.0), y, Inches(0.4), Inches(0.4),
                 f'{i+1}.', size=12, bold=True, color=ACCENT)
        add_text(s, Inches(7.4), y, Inches(5.8), Inches(0.4),
                 step, size=12, color=FG)

slide_app_store_submit()


# ───────────────────────────────────────────────────────────
# PART 5. 체크리스트 & 문제해결
# ───────────────────────────────────────────────────────────
def slide_checklist():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, 21, 22, 'PART 5')
    add_title(s, '배포 전 체크리스트', '꼭 확인할 것들')

    columns = [
        ('✅ 코드', [
            'viewport meta 설정',
            'font-size 16px 이상',
            '터치 영역 44px 이상',
            'touch-action: none 캔버스',
            'safe-area-inset 적용',
            'Service Worker 등록',
            'manifest.json 연결',
            '가로/세로 회전 대응',
        ]),
        ('🧪 테스트', [
            '시뮬레이터 iPhone SE/Pro Max',
            '실기기 터치 반응 확인',
            '3D 드래그 vs OrbitControls',
            '뒤로가기 버튼 동작',
            '되돌리기 20단계 확인',
            'JSON 저장/불러오기',
            '스크린샷 저장',
            'PWA 오프라인 동작',
        ]),
        ('📱 배포', [
            'Bundle ID 유일성',
            '아이콘 1024×1024 준비',
            '스크린샷 모든 사이즈',
            '프라이버시 정책 URL',
            '앱 설명 (한/영)',
            '테스트 계정 정보',
            '심사 가이드라인 검토',
            'TestFlight 내부 테스트',
        ]),
    ]
    for i, (title, items) in enumerate(columns):
        x = Inches(0.6 + i * 4.25)
        card = add_rect(s, x, Inches(2.0), Inches(4.1), Inches(5.2))
        add_text(s, x + Inches(0.25), Inches(2.15), Inches(3.7), Inches(0.5),
                 title, size=16, bold=True, color=ACCENT)
        for j, it in enumerate(items):
            y = Inches(2.75 + j * 0.5)
            # 체크박스
            cb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.25), y + Inches(0.08), Inches(0.25), Inches(0.25))
            cb.fill.solid(); cb.fill.fore_color.rgb = BG
            cb.line.color.rgb = ACCENT; cb.line.width = Pt(1)
            add_text(s, x + Inches(0.6), y, Inches(3.4), Inches(0.4),
                     it, size=12, color=FG)

slide_checklist()


def slide_troubleshoot():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, 22, 22, 'PART 5')
    add_title(s, '자주 겪는 이슈 & 해결', 'Troubleshooting')

    issues = [
        ('🔴 3D 터치 드래그가 카메라 회전과 충돌',
         'pointerdown에서 raycaster로 객체 체크 → 히트 시 orbit.enabled = false. m1.html에 이미 구현됨.'),
        ('🔴 iPhone 입력 시 화면이 줌됨',
         '<meta name="viewport" ...maximum-scale=1, user-scalable=no>. input 폰트 16px 이상.'),
        ('🔴 노치 영역에 UI 가려짐',
         'viewport-fit=cover + padding에 env(safe-area-inset-top/bottom) 사용. 이미 적용됨.'),
        ('🔴 Service Worker가 캐시한 구버전이 계속 뜸',
         'sw.js의 CACHE 이름을 "labplace-m1-v2" 처럼 버전 올리고 재배포.'),
        ('🔴 Xcode 빌드 실패: "No such module Capacitor"',
         'cd ios/App && pod install 실행. 그 후 .xcworkspace 파일로 열기 (.xcodeproj 아님).'),
        ('🔴 App Store 심사 거절 - "웹뷰만 있는 앱"',
         '가이드라인 4.2: 오프라인 기능, 네이티브 통합(저장소/공유 등), 고유 가치를 명시. 실제로 동작.'),
        ('🔴 TestFlight 빌드 90일 만료',
         '새 빌드 재업로드. 빌드 번호만 올리면 됨. 버전은 유지 가능.'),
        ('🔴 Three.js 로딩 실패 (오프라인)',
         'sw.js에서 unpkg.com을 캐시에 포함. 또는 번들러로 로컬 인라인.'),
    ]
    for i, (q, a) in enumerate(issues):
        y = Inches(1.95 + i * 0.63)
        add_text(s, Inches(0.6), y, Inches(12.2), Inches(0.3),
                 q, size=12, bold=True, color=WARN)
        add_text(s, Inches(0.8), y + Inches(0.28), Inches(12), Inches(0.35),
                 '→ ' + a, size=11, color=FG_DIM)

slide_troubleshoot()


# ───────────────────────────────────────────────────────────
# 저장
# ───────────────────────────────────────────────────────────
out = 'LabPlace_m1_배포가이드.pptx'
prs.save(out)
print(f'✅ PPT 생성 완료: {out}')
print(f'   총 슬라이드: {len(prs.slides)}장')
